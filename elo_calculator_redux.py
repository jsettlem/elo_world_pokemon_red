import random
from typing import Iterable

import numpy as np

from common_methods import *
import numpy
from numpy import array, linspace
from sklearn.neighbors import KernelDensity
from matplotlib import pyplot
from scipy.signal import argrelextrema
from sklearn import linear_model
from hash_id_common import display_hashid

DO_POWERPOINT = True


def greatest_diffs_cluster(elo_list, cluster_count=14):
	boundaries = sorted(list(range(len(elo_list) - 1)), key=lambda i: elo_list[i + 1] - elo_list[i], reverse=True)
	return sorted([elo_list[b] for b in boundaries[:cluster_count]])


def kde_cluster(elo_list: List[int], K=10, save=False, index=0):
	pyplot.clf()
	min_elo, max_elo = min(elo_list), max(elo_list)

	elo_array = array(elo_list).reshape(-1, 1)
	kde = KernelDensity(bandwidth=K).fit(elo_array)

	space = linspace(min_elo, max_elo, num=int(max_elo - min_elo) * 2)
	estimates = kde.score_samples(space.reshape(-1, 1))

	mins = argrelextrema(estimates, numpy.less)[0]

	pyplot.plot(space, estimates)
	pyplot.plot(space[mins], estimates[mins], 'ro')
	pyplot.title(f"K: {K:.1f} - Number of minima: {len(mins)}")

	if save:
		pyplot.savefig(f"kdes/{index}.png", dpi=320)

	return [space[k] for k in mins]


def update_elo(t1: Trainer, t2: Trainer, winner: str, K: int = 32):
	r1 = 10 ** (t1.elo / 400)
	r2 = 10 ** (t2.elo / 400)

	e1 = r1 / (r1 + r2)
	e2 = r2 / (r1 + r2)

	if winner == "trainer":
		s1, s2 = 1, 0
	elif winner == "enemy":
		s1, s2 = 0, 1
	else:
		s1, s2 = 0.5, 0.5

	t1.elo = t1.elo + K * (s1 - e1)
	t2.elo = t2.elo + K * (s2 - e2)


def battles_with_wasted_heal(battles: Iterable[Battle]):
	for battle in battles:
		if "movie" not in battle.original_path and "POTION" in battle.turns[0].action.item:
			print(battle.original_path)

def generate_lr_elo(battles: List[Battle], trainers: List[Trainer]):
	N = len(trainers)
	X = []
	Y = []
	for battle in battles:
		t1 = trainers.index(battle.player)
		t2 = trainers.index(battle.enemy)
		v = np.zeros(N)
		v[t1] = 1
		v[t2] = -1
		X.append(v)
		if battle.winner == "Draw":
			Y.append(1)
			Y.append(0)
			X.append(v)
		else:
			Y.append(battle.winner == "trainer")

	clf = linear_model.LogisticRegression()
	clf.fit(X, Y)
	return list(clf.coef_[0] * 173 + 1500), clf.intercept_[0]


def main():
	trainer_dict, battle_dict = load_pickle("omega_redux.pickle")

	trainer_wins = enemy_wins = draws = 0
	trainer_list = list(trainer_dict.values())
	battle_list = list(battle_dict.values())

	for battle in battle_list:
		if battle.winner == 'trainer':
			trainer_wins += 1
		elif battle.winner == 'enemy':
			enemy_wins += 1
		else:
			draws += 1
	for _ in range(1):
		for K in (32,):# 16, 8, 4, 2, 1, 0.5, 0.25):
			random.shuffle(battle_list)
			for battle in battle_list:
				if battle.player != battle.enemy:
					update_elo(battle.player, battle.enemy, battle.winner, K=K)
		battle_trainers = list(trainer_dict.values())
		print([t.elo for t in battle_trainers])

		# battle_trainers.sort(key=lambda t: t.elo)
		# battles_trainers_sorted = battle_trainers.copy()
		# battles_trainers_sorted.sort(key=lambda t: t.class_id * 1000 + t.instance_id)
		# order = [battle_trainers.index(trainer) for trainer in battles_trainers_sorted]
		# print(order)

	regression_elo = generate_lr_elo(battle_list, trainer_list)
	for i, trainer in enumerate(trainer_list):
		trainer.lr_elo = regression_elo[0][i]

	trainer_list.sort(key=lambda t: abs(t.lr_elo - t.elo))
	print([t.elo - t.lr_elo for t in trainer_list])

	print("Trainer wins:", trainer_wins, "enemy wins:", enemy_wins, "draws:", draws)

	trainer_list.sort(key=lambda t: t.lr_elo)
	elo_list = [t.lr_elo for t in trainer_list]
	kde_boundaries = kde_cluster(elo_list, K=28)


	# for b in range(1, 281):
	# 	kde_cluster(elo_list, K=b/10, save=True, index=b)

	greatest_diff_boundaries = greatest_diffs_cluster(elo_list, cluster_count=14)
	human_boundaries = [610.52, 738.02, 1022.02, 1127.38, 1290.28,
	                    1408.39, 1485.61, 1543.89, 1655.25, 1808.88, 1887.05,
	                    1994.33, 2154.91, 2546.942738]

	print(greatest_diff_boundaries)
	current_kde_tier = 0
	current_diff_tier = 0
	current_human_tier = 0
	for trainer in trainer_list:
		if current_kde_tier < len(kde_boundaries) and trainer.lr_elo > kde_boundaries[current_kde_tier]:
			current_kde_tier += 1
		if current_diff_tier < len(greatest_diff_boundaries) and trainer.lr_elo > greatest_diff_boundaries[
			current_diff_tier]:
			current_diff_tier += 1
		if current_human_tier < len(human_boundaries) and trainer.lr_elo >= human_boundaries[current_human_tier] - 0.1:
			current_human_tier += 1
		trainer.tier = current_human_tier
		win_count, lose_count, draw_count = trainer.get_win_rate()
		print("\t".join(
			("trainer:", trainer.identifier, "location", trainer.location if trainer.location else "somewhere",
			 "party:", ", ".join(f"Lv. {mon.level} {mon.species}" for mon in trainer.party_mons),
			 "class:", str(trainer.class_id), "instance:", str(trainer.instance_id),
			 "win count:", str(win_count), "lose count:", str(lose_count),
			 "draw count:", str(draw_count), "LR elo:", str(trainer.lr_elo), "elo:", str(trainer.elo),
			 "kde tier:", str(current_kde_tier), "diff tier:", str(current_diff_tier))))

	battle_list.sort(key=lambda b: b.losing_trainer.lr_elo - b.winning_trainer.lr_elo, reverse=True)

	lorelei_draws = [b for b in battle_list if b.player.class_id == 244 and b.winner.startswith("Draw")]
	print("")
	# for upset in battle_list[:100]:
	# 	print("\t".join((upset.original_path, upset.winning_trainer.identifier, str(upset.winning_trainer.lr_elo),
	# 	                 upset.losing_trainer.identifier, str(upset.losing_trainer.lr_elo),
	# 	                 str(upset.losing_trainer.lr_elo - upset.winning_trainer.lr_elo))))

	# battle_list.sort(key=lambda b: len(b.turns), reverse=True)
	# for long_battle in battle_list[:100]:
	# 	print("\t".join(str(s) for s in (long_battle.seed, long_battle.winner, long_battle.winning_trainer.identifier,
	# 	                long_battle.losing_trainer.identifier, len(long_battle.turns))))

	if DO_POWERPOINT:
		orig_trainer_dict, orig_battle_dict = load_pickle("../omega.pickle")
		orig_battle_list = list(orig_battle_dict.values())
		for battle in orig_battle_list:
			if battle.player != battle.enemy:
				update_elo(battle.player, battle.enemy, battle.winner)

		orig_battle_trainers = list(orig_trainer_dict.values())
		orig_battle_trainers.sort(key=lambda t: t.elo)

		trainer_list = list(trainer_dict.values())
		trainer_list.sort(key=lambda t: t.lr_elo)

		for trainer in trainer_list:
			decisive_battles = [b for b in trainer.battles if "Draw" not in b.winner]
			victories = [b for b in decisive_battles if b.winning_trainer == trainer]
			defeats = [b for b in decisive_battles if b.losing_trainer == trainer]
			trainer.greatest_victory = max(victories, key=lambda b: b.losing_trainer.lr_elo)
			trainer.greatest_defeat = min(defeats, key=lambda b: b.winning_trainer.lr_elo)
		tier_names = ["F", "D-", "D", "D+", "C-", "C",
		              "C+", "B-", "B", "B+", "A-", "A", "A+",
		              "S", "S+"]
		tier_colors = ["00b050", "24bb45", "4bc735", "6fd226", "93de15", "b8e900",
		               "dcf400", "ffff00", "ffd600", "ffac00", "ff5d00", "ff5700",
		               "ff2b00", "ff0000", "ff0050"]

		from pptx import Presentation
		from pptx.dml.color import RGBColor

		prs = Presentation("../video-redux/trainer-cards/template.pptm")
		# prs.slide_height = prs.slide_height // 4
		# prs.slide_width = prs.slide_width // 4
		lose_width = 4254918
		pokemon_in_order = ["Bulbasaur", "Ivysaur", "Venusaur", "Charmander", "Charmeleon", "Charizard", "Squirtle",
		                    "Wartortle", "Blastoise", "Caterpie", "Metapod", "Butterfree", "Weedle", "Kakuna",
		                    "Beedrill", "Pidgey", "Pidgeotto", "Pidgeot", "Rattata", "Raticate", "Spearow", "Fearow",
		                    "Ekans", "Arbok", "Pikachu", "Raichu", "Sandshrew", "Sandslash", "Nidoran♀", "Nidorina",
		                    "Nidoqueen", "Nidoran♂", "Nidorino", "Nidoking", "Clefairy", "Clefable", "Vulpix",
		                    "Ninetales", "Jigglypuff", "Wigglytuff", "Zubat", "Golbat", "Oddish", "Gloom", "Vileplume",
		                    "Paras", "Parasect", "Venonat", "Venomoth", "Diglett", "Dugtrio", "Meowth", "Persian",
		                    "Psyduck", "Golduck", "Mankey", "Primeape", "Growlithe", "Arcanine", "Poliwag", "Poliwhirl",
		                    "Poliwrath", "Abra", "Kadabra", "Alakazam", "Machop", "Machoke", "Machamp", "Bellsprout",
		                    "Weepinbell", "Victreebel", "Tentacool", "Tentacruel", "Geodude", "Graveler", "Golem",
		                    "Ponyta", "Rapidash", "Slowpoke", "Slowbro", "Magnemite", "Magneton", "Farfetch'D", "Doduo",
		                    "Dodrio", "Seel", "Dewgong", "Grimer", "Muk", "Shellder", "Cloyster", "Gastly", "Haunter",
		                    "Gengar", "Onix", "Drowzee", "Hypno", "Krabby", "Kingler", "Voltorb", "Electrode",
		                    "Exeggcute", "Exeggutor", "Cubone", "Marowak", "Hitmonlee", "Hitmonchan", "Lickitung",
		                    "Koffing", "Weezing", "Rhyhorn", "Rhydon", "Chansey", "Tangela", "Kangaskhan", "Horsea",
		                    "Seadra", "Goldeen", "Seaking", "Staryu", "Starmie", "Mr.Mime", "Scyther", "Jynx",
		                    "Electabuzz", "Magmar", "Pinsir", "Tauros", "Magikarp", "Gyarados", "Lapras", "Ditto",
		                    "Eevee", "Vaporeon", "Jolteon", "Flareon", "Porygon", "Omanyte", "Omastar", "Kabuto",
		                    "Kabutops", "Aerodactyl", "Snorlax", "Articuno", "Zapdos", "Moltres", "Dratini",
		                    "Dragonair", "Dragonite", "Mewtwo", "Mew"]
		for i, slide in enumerate(prs.slides):
			trainer = trainer_list[i]
			battle_count = trainer.win_count + trainer.draw_count + trainer.lose_count

			trainer_rank = 391 - i
			trainer_old_rank = 391 - next(i for i,t in enumerate(orig_battle_trainers) if t.identifier == trainer.identifier)
			rank_diff = trainer_old_rank - trainer_rank
			for shape in slide.shapes:
				if shape.name == "win":
					shape.width = int(trainer.win_count / battle_count * lose_width)
				elif shape.name == "draw":
					shape.width = int((trainer.win_count + trainer.draw_count) / battle_count * lose_width)
				elif shape.name == "tier":
					shape.fill.fore_color.rgb = RGBColor.from_string(tier_colors[trainer.tier])
				elif shape.name == "pic":
					image, rid = shape.part.get_or_add_image_part(
						"V:\\Dropbox\\elo-world\\video\\3d\\textures\\trainer_sprites\\" +
						get_trainer_by_id(trainer.class_id, trainer.instance_id)[0]["sprite"])
					shape._element.blipFill.blip.rEmbed = rid
					shape.shadow._element.getchildren()[3].getchildren()[0].getchildren()[0].attrib["val"] = tier_colors[trainer.tier]
				elif shape.name == "lose":
					pass
				elif shape.name == "wlpic":
					image, rid = shape.part.get_or_add_image_part(
						"V:\\Dropbox\\elo-world\\video\\3d\\textures\\trainer_sprites\\" +
						get_trainer_by_id(trainer.greatest_defeat.winning_trainer.class_id, trainer.greatest_defeat.winning_trainer.instance_id)[0]["sprite"])
					shape._element.blipFill.blip.rEmbed = rid
					shape.shadow._element.getchildren()[3].getchildren()[0].getchildren()[0].attrib["val"] = tier_colors[trainer.greatest_defeat.winning_trainer.tier]
				elif shape.name == "gwpic":
					image, rid = shape.part.get_or_add_image_part(
						"V:\\Dropbox\\elo-world\\video\\3d\\textures\\trainer_sprites\\" +
						get_trainer_by_id(trainer.greatest_victory.losing_trainer.class_id, trainer.greatest_victory.losing_trainer.instance_id)[0]["sprite"])
					shape._element.blipFill.blip.rEmbed = rid
					shape.shadow._element.getchildren()[3].getchildren()[0].getchildren()[0].attrib["val"] = tier_colors[trainer.greatest_victory.losing_trainer.tier]
				elif shape.name.startswith("pk"):
					pk_index = int(shape.name[2]) - 1
					if pk_index < len(trainer.party_mons):
						species = trainer.party_mons[pk_index].species
						dex_number = pokemon_in_order.index(species.title()) + 1
						image, rid = shape.part.get_or_add_image_part(
							f"Z:\\Dropbox\\elo-world\\video-redux\\trainer-cards\\pokemon-sprites\\scaled\\{dex_number:03}MS6.png")
						shape._element.blipFill.blip.rEmbed = rid
					else:
						image, rid = shape.part.get_or_add_image_part(
							"Z:\\Dropbox\\elo-world\\video-redux\\trainer-cards\\pokemon-sprites\\scaled\\XYLoadingMS-2.png")
						shape._element.blipFill.blip.rEmbed = rid
				elif shape.name == "up-trianggle":
					if rank_diff < 0:
						shape.fill.background()
						shape.line.fill.background()
				elif shape.name == "down-triangle":
					if rank_diff >= 0:
						shape.fill.background()
						shape.line.fill.background()
				if shape.has_text_frame:
					for p in shape.text_frame.paragraphs:
						for r in p.runs:
							if r.text == "t":
								r.text = " " + tier_names[trainer.tier]
							elif r.text == "ra":
								r.text = str(len(trainer_list) - i)
							elif r.text == "na":
								r.text = trainer.name
							elif r.text == "nu":
								r.text = str(trainer.instance_id)
							elif r.text == "wld":
								r.text = f"W-D-L:{trainer.win_count}-{trainer.draw_count}-{trainer.lose_count}"
							elif r.text == "elo":
								r.text = str(int(trainer.lr_elo))
							elif r.text == "loc":
								r.text = trainer.location if trainer.location else "somewhere"
							elif r.text == "dtr":
								r.text = tier_names[trainer.greatest_defeat.winning_trainer.tier]
								r.font.color.rgb = RGBColor.from_string(tier_colors[trainer.greatest_defeat.winning_trainer.tier])
							elif r.text == "dna":
								r.text = trainer.greatest_defeat.winning_trainer.name
							elif r.text == "dnu":
								r.text = f"{trainer.greatest_defeat.winning_trainer.instance_id} ({int(trainer.greatest_defeat.winning_trainer.lr_elo)})"
							elif r.text == "dh":
								r.text = display_hashid(trainer.greatest_defeat.seed)
							elif r.text == "vtr":
								r.text = tier_names[trainer.greatest_victory.losing_trainer.tier]
								r.font.color.rgb = RGBColor.from_string(
									tier_colors[trainer.greatest_victory.losing_trainer.tier])
							elif r.text == "vna":
								r.text = trainer.greatest_victory.losing_trainer.name
							elif r.text == "vnu":
								r.text = f"{trainer.greatest_victory.losing_trainer.instance_id} ({int(trainer.greatest_victory.losing_trainer.lr_elo)})"
							elif r.text == "vh":
								r.text = display_hashid(trainer.greatest_victory.seed)
							elif r.text == "old":
								r.text = f"Old avg: Lv. {sum(p.level for p in trainer.party_mons) // len(trainer.party_mons)}"
							elif r.text == " gai":
								r.text = " " + str(rank_diff) if rank_diff >= 0 else " "
							elif r.text == " los":
								r.text = " " + str(abs(rank_diff)) if rank_diff < 0 else " "
		prs.save("test1_redux.pptm")


if __name__ == '__main__':
	main()
