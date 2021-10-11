from typing import Iterable



from common_methods import *
import numpy
from numpy import array, linspace
from sklearn.neighbors import KernelDensity
from matplotlib import pyplot
from scipy.signal import argrelextrema

DO_POWERPOINT = False


def greatest_diffs_cluster(elo_list, cluster_count=14):
	boundaries = sorted(list(range(len(elo_list) - 1)), key=lambda i: elo_list[i + 1] - elo_list[i], reverse=True)
	return sorted([elo_list[b] for b in boundaries[:cluster_count]])


def kde_cluster(elo_list: List[int], K=10, save=False, index=0):
	pyplot.clf()
	min_elo, max_elo = min(elo_list), max(elo_list)

	elo_array = array(elo_list).reshape(-1, 1)
	kde = KernelDensity(bandwidth=K).fit(elo_array)

	space = linspace(min_elo, max_elo, num=int(max_elo - min_elo) * 2)
	estimates = kde.score_samples(space.reshape(-1, 1))

	mins = argrelextrema(estimates, numpy.less)[0]

	pyplot.plot(space, estimates)
	pyplot.plot(space[mins], estimates[mins], 'ro')
	pyplot.title(f"K: {K:.1f} - Number of minima: {len(mins)}")

	if save:
		pyplot.savefig(f"kdes/{index}.png", dpi=320)

	return [space[k] for k in mins]


def update_elo(t1: Trainer, t2: Trainer, winner: str, K: int = 32):
	r1 = 10 ** (t1.elo / 400)
	r2 = 10 ** (t2.elo / 400)

	e1 = r1 / (r1 + r2)
	e2 = r2 / (r1 + r2)

	if winner == "trainer":
		s1, s2 = 1, 0
	elif winner == "enemy":
		s1, s2 = 0, 1
	else:
		s1, s2 = 0.5, 0.5

	t1.elo = t1.elo + K * (s1 - e1)
	t2.elo = t2.elo + K * (s2 - e2)


def battles_with_wasted_heal(battles: Iterable[Battle]):
	for battle in battles:
		if "movie" not in battle.original_path and "POTION" in battle.turns[0].action.item:
			print(battle.original_path)


def main():
	trainer_dict, battle_dict = load_pickle("../omega.pickle")

	trainer_wins = enemy_wins = draws = 0
	battle_list = list(battle_dict.values())

	lorelei_draws = [b for b in battle_list if b.player.class_id == 244 and b.winner.startswith("Draw")]
	pokemon_lorelei_drew_against = [p for b in lorelei_draws for p in b.enemy.party_mons]
	print("Pokémon Lorelei drew against:")
	print("\n".join([f"Lv. {p.level} {p.species}" for p in pokemon_lorelei_drew_against]))
	
	bug_draws = [b for b in battle_list if "bgb" in b.winner]
	print(len(bug_draws))
	karate_master_wins = [b for b in battle_list if b.winning_trainer.class_id == 224 and b.winning_trainer.instance_id == 1 and b.losing_trainer.class_id == 239]
	lance_wins_against_oak = [b for b in battle_list if b.winning_trainer.class_id == 247 and b.losing_trainer.class_id == 226 and b.losing_trainer.instance_id == 1]

	battles_with_wasted_heal(battle_list)

	for battle in battle_list:
		if battle.winner == 'trainer':
			trainer_wins += 1
		elif battle.winner == 'enemy':
			enemy_wins += 1
		else:
			draws += 1
		if battle.player != battle.enemy:
			update_elo(battle.player, battle.enemy, battle.winner)

	oddish_battles = [b for b in battle_list if b.enemy.party_mons[0].species == "ODDISH" and b.player.party_mons[0].species == "ODDISH" and "Draw" in b.winner]

	print("Trainer wins:", trainer_wins, "enemy wins:", enemy_wins, "draws:", draws)
	battle_trainers = list(trainer_dict.values())
	battle_trainers.sort(key=lambda t: t.elo)

	elo_list = [t.elo for t in battle_trainers]
	kde_boundaries = kde_cluster(elo_list, K=28)

	for b in range(1, 281):
		kde_cluster(elo_list, K=b/10, save=True, index=b)

	greatest_diff_boundaries = greatest_diffs_cluster(elo_list, cluster_count=14)
	human_boundaries = [444.6391, 561.3502, 1038.7349, 1118.5088, 1225.8032,
	                    1358.1477, 1500.3403, 1661.6945, 1782.0818, 1887.0274, 2214.7272,
	                    2370.2143, 2559.3289]

	print(greatest_diff_boundaries)
	current_kde_tier = 0
	current_diff_tier = 0
	current_human_tier = 0
	for trainer in battle_trainers:
		if current_kde_tier < len(kde_boundaries) and trainer.elo > kde_boundaries[current_kde_tier]:
			current_kde_tier += 1
		if current_diff_tier < len(greatest_diff_boundaries) and trainer.elo > greatest_diff_boundaries[
			current_diff_tier]:
			current_diff_tier += 1
		if current_human_tier < len(human_boundaries) and trainer.elo >= human_boundaries[current_human_tier] - 0.1:
			current_human_tier += 1
		trainer.tier = current_human_tier
		win_count, lose_count, draw_count = trainer.get_win_rate()
		print("\t".join(
			("trainer:", trainer.identifier, "location", trainer.location if trainer.location else "somewhere",
			 "party:", ", ".join(f"Lv. {mon.level} {mon.species}" for mon in trainer.party_mons),
			 "class:", str(trainer.class_id), "instance:", str(trainer.instance_id),
			 "win count:", str(win_count), "lose count:", str(lose_count),
			 "draw count:", str(draw_count), "elo:", str(trainer.elo),
			 "kde tier:", str(current_kde_tier), "diff tier:", str(current_diff_tier))))

	battle_list.sort(key=lambda b: b.losing_trainer.elo - b.winning_trainer.elo, reverse=True)
	for upset in battle_list[:100]:
		print("\t".join((upset.original_path, upset.winning_trainer.identifier, str(upset.winning_trainer.elo),
		                 upset.losing_trainer.identifier, str(upset.losing_trainer.elo),
		                 str(upset.losing_trainer.elo - upset.winning_trainer.elo))))

	if DO_POWERPOINT:
		trainer_list = list(trainer_dict.values())
		trainer_list.sort(key=lambda t: t.elo)

		for trainer in trainer_list:
			decisive_battles = [b for b in trainer.battles if "Draw" not in b.winner]
			victories = [b for b in decisive_battles if b.winning_trainer == trainer]
			defeats = [b for b in decisive_battles if b.losing_trainer == trainer]
			trainer.greatest_victory = max(victories, key=lambda b: b.losing_trainer.elo)
			trainer.greatest_defeat = min(defeats, key=lambda b: b.winning_trainer.elo)

		tier_names = ["F", "D-", "D", "D+", "C-", "C",
		              "C+", "B-", "B", "B+", "A-", "A", "A+",
		              "S", "S+"]
		tier_colors = ["00b050", "24bb45", "4bc735", "6fd226", "93de15", "b8e900",
		               "dcf400", "ffff00", "ffd600", "ffac00", "ff5d00", "ff5700",
		               "ff2b00", "ff0000", "ff0000"]

		lose_width = 3741321

		from pptx import Presentation
		from pptx.dml.color import RGBColor

		prs = Presentation("base.pptm")
		prs.slide_height = prs.slide_height // 4
		for i, slide in enumerate(prs.slides):
			trainer = trainer_list[i]
			battle_count = trainer.win_count + trainer.draw_count + trainer.lose_count
			for shape in slide.shapes:
				if shape.name == "win":
					shape.width = int(trainer.win_count / battle_count * lose_width)
				elif shape.name == "draw":
					shape.width = int((trainer.win_count + trainer.draw_count) / battle_count * lose_width)
				elif shape.name == "tier":
					shape.fill.fore_color.rgb = RGBColor.from_string(tier_colors[trainer.tier])
				elif shape.name == "pic":
					image, rid = shape.part.get_or_add_image_part(
						"V:\\Dropbox\\elo-world\\video\\3d\\textures\\trainer_sprites\\" +
						get_trainer_by_id(trainer.class_id, trainer.instance_id)[0]["sprite"])
					shape._element.blipFill.blip.rEmbed = rid
				elif shape.name == "pkmn":
					for i, cell in enumerate(shape.table.iter_cells()):
						for p in cell.text_frame.paragraphs:
							for r in p.runs:
								if i < len(trainer.party_mons):
									if "Lv" in r.text:
										r.text = f"Lv. {trainer.party_mons[i].level}"
									else:
										r.text = f"{trainer.party_mons[i].species}"
								else:
									r.text = " "
				elif shape.name == "upset":
					for cell in shape.table.iter_cells():
						for p in cell.text_frame.paragraphs:
							for r in p.runs:
								if r.text == "dna":
									r.text = f"{tier_names[trainer.greatest_defeat.winning_trainer.tier]} {trainer.greatest_defeat.winning_trainer.name}"
								elif r.text == "dnu":
									r.text = f"{trainer.greatest_defeat.winning_trainer.instance_id} - {int(trainer.greatest_defeat.winning_trainer.elo)}"
								elif r.text == "vna":
									r.text = f"{tier_names[trainer.greatest_victory.losing_trainer.tier]} {trainer.greatest_victory.losing_trainer.name}"
								elif r.text == "vnu":
									r.text = f"{trainer.greatest_victory.losing_trainer.instance_id} - {int(trainer.greatest_victory.losing_trainer.elo)}"

				if shape.has_text_frame:
					for p in shape.text_frame.paragraphs:
						for r in p.runs:
							if r.text == "t":
								r.text = tier_names[trainer.tier]
							elif r.text == "ra":
								r.text = str(len(trainer_list) - i)
							elif r.text == "na":
								r.text = trainer.name
							elif r.text == "nu":
								r.text = str(trainer.instance_id)
							elif r.text == "wld":
								r.text = f"W-D-L:{trainer.win_count}-{trainer.draw_count}-{trainer.lose_count}"
							elif r.text == "elo":
								r.text = str(int(trainer.elo))
							elif r.text == "loc":
								r.text = trainer.location

		prs.save("test1.pptm")


if __name__ == '__main__':
	main()
